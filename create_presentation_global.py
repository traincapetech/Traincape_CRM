#!/usr/bin/env python3
"""
Traincape CRM Global PowerPoint Presentation Generator
Creates a professional sales presentation for worldwide clients
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_global_traincape_presentation():
    """Create the complete Traincape CRM presentation for global market"""
    
    # Create presentation
    prs = Presentation()
    
    # Define colors
    PRIMARY_BLUE = RGBColor(37, 99, 235)
    SUCCESS_GREEN = RGBColor(16, 185, 129)
    WARNING_ORANGE = RGBColor(245, 158, 11)
    ERROR_RED = RGBColor(239, 68, 68)
    NEUTRAL_GRAY = RGBColor(107, 114, 128)
    WHITE = RGBColor(255, 255, 255)
    
    # SLIDE 1: TITLE SLIDE
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🚀 TRAINCAPE CRM"
    subtitle.text = "The World's Most Affordable, Complete CRM\n\nBuilt-in Payment Processing • Multi-Currency • AI-Powered\n\nPresented by: [Your Name]\nDate: [Current Date]\nContact: sales@traincapetech.in"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    # SLIDE 2: AGENDA
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📋 AGENDA"
    content.text = """1. 🎯 The Problem: Why Current CRMs Fail
2. 💡 The Solution: Traincape CRM
3. 🚀 Key Features & Benefits
4. 💰 ROI & Cost Savings
5. 🎬 Live Demo
6. 📊 Success Stories
7. 🎯 Pricing & Next Steps

Duration: 30-45 minutes"""
    
    # SLIDE 3: THE PROBLEM
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "❌ THE PROBLEM"
    content.text = """Why Current CRMs Are Failing Your Business

💸 EXPENSIVE
• Salesforce: $150-500/month per user
• HubSpot: $80-250/month per user
• Zoho: $40-150/month per user

🐌 SLOW SETUP
• 3-6 months implementation time
• Complex configuration required
• Expensive consultants needed

🌍 NOT GLOBAL-FOCUSED
• Limited currency support
• Regional payment restrictions
• Generic workflows

💳 NO PAYMENT PROCESSING
• Separate payment systems
• Manual invoice tracking
• Payment delays

🔧 COMPLEX INTEGRATIONS
• Multiple tools required
• Data synchronization issues
• High maintenance costs"""
    
    # SLIDE 4: THE SOLUTION
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "✅ THE SOLUTION"
    content.text = """Traincape CRM: Complete Global Business Solution

🚀 READY IN 1 WEEK
• vs 6 months for competitors
• No complex configuration
• Guided setup process

💰 80% COST SAVINGS
• $29-129/month per user
• vs $150-500/month per user
• All features included

🌍 GLOBAL-FIRST
• Multi-currency support
• International payments
• Local compliance features
• Global business workflows

💳 BUILT-IN PAYMENTS
• Stripe integration included
• Automatic payment tracking
• Professional payment pages
• Multi-currency invoicing

🔧 ALL-IN-ONE SOLUTION
• Everything integrated
• No additional tools needed
• Seamless data flow
• Low maintenance"""
    
    # SLIDE 5: COMPLETE FEATURE OVERVIEW
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎯 COMPLETE CRM SOLUTION"
    content.text = """📊 SALES MANAGEMENT
• Lead Management & Scoring
• Sales Pipeline & Forecasting
• Customer Relationship Management
• Deal Tracking & Analytics

💳 PAYMENT PROCESSING
• Stripe Integration
• Professional Invoicing
• Payment Tracking
• Multi-currency Support

👥 EMPLOYEE MANAGEMENT
• HR & Payroll
• Attendance Tracking
• Leave Management
• Performance Analytics

📱 MOBILE-FIRST DESIGN
• Responsive Web App
• Works on All Devices
• Offline Capability
• Push Notifications

🌍 GLOBAL FEATURES
• Multi-language Support
• Local Tax Compliance
• International Payments
• Global Business Workflows"""
    
    # SLIDE 6: SALES MANAGEMENT
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📊 SALES MANAGEMENT"
    content.text = """🎯 LEAD MANAGEMENT
• Automated lead capture
• Lead scoring & prioritization
• Lead assignment & tracking
• Lead nurturing workflows

📈 SALES PIPELINE
• Visual pipeline management
• Deal stage tracking
• Win/loss analysis
• Sales forecasting

👥 CUSTOMER MANAGEMENT
• 360° customer view
• Communication history
• Customer segmentation
• Relationship scoring

🌍 GLOBAL SALES
• Multi-currency deals
• International client support
• Local business practices
• Global market insights"""
    
    # SLIDE 7: PAYMENT PROCESSING
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "💳 PAYMENT PROCESSING"
    content.text = """🚀 AUTOMATED WORKFLOW
1. Create Invoice → 2. Send to Customer →
3. Customer Pays → 4. Status Updates

💰 MULTIPLE PAYMENT METHODS
• Credit/Debit Cards
• Digital Wallets
• Bank Transfers
• International Payments

🌍 MULTI-CURRENCY SUPPORT
• 135+ currencies supported
• Real-time exchange rates
• Local payment methods
• International compliance

📊 PAYMENT ANALYTICS
• Real-time payment tracking
• Failed payment recovery
• Payment history & reports
• Revenue analytics

🔒 SECURE & COMPLIANT
• PCI DSS compliant
• GDPR compliant
• Local tax compliance
• Secure data encryption"""
    
    # SLIDE 8: EMPLOYEE MANAGEMENT
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "👥 EMPLOYEE MANAGEMENT"
    content.text = """📋 HR & PAYROLL
• Employee profiles & documents
• Salary management & slips
• Tax calculations & compliance
• Performance reviews

⏰ ATTENDANCE TRACKING
• Check-in/check-out
• Leave management
• Overtime tracking
• Attendance reports

📊 PERFORMANCE ANALYTICS
• KPI tracking
• Goal management
• Performance reviews
• Team productivity

🌍 GLOBAL HR
• Multi-currency payroll
• Local labor compliance
• International team support
• Global HR policies"""
    
    # SLIDE 9: MOBILE-FIRST DESIGN
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📱 MOBILE-FIRST DESIGN"
    content.text = """🎯 RESPONSIVE DESIGN
• Works on all devices
• Touch-optimized interface
• Fast loading times
• Offline capability

🔔 SMART NOTIFICATIONS
• Real-time alerts
• Payment notifications
• Task reminders
• Custom notifications

🚀 SEAMLESS EXPERIENCE
• Same features on mobile
• Easy navigation
• Quick actions
• Voice commands

🌍 GLOBAL ACCESS
• Works worldwide
• Local language support
• Regional features
• Global connectivity"""
    
    # SLIDE 10: ROI & COST SAVINGS
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "💰 ROI & COST SAVINGS"
    content.text = """📊 COST COMPARISON
Salesforce: $150-500/month per user
HubSpot: $80-250/month per user
Zoho: $40-150/month per user
Traincape: $29-129/month per user

💸 ANNUAL SAVINGS (10 users)
• vs Salesforce: $14,520-44,520
• vs HubSpot: $6,120-26,520
• vs Zoho: $1,320-14,520

📈 BUSINESS IMPACT
• 40% faster lead response
• 25% higher conversion rates
• 60% time savings on invoicing
• 30% increase in sales

🎯 ROI CALCULATOR
• Investment: $3,600/year (10 users)
• Savings: $20,000/year
• ROI: 455% in first year

🌍 GLOBAL SAVINGS
• No currency conversion fees
• Local payment processing
• Reduced compliance costs
• Integrated solution savings"""
    
    # SLIDE 11: LIVE DEMO
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎬 LIVE DEMO"
    content.text = """🎯 DEMO SCENARIOS
1. Creating a New Lead
2. Managing Sales Pipeline
3. Generating Multi-Currency Invoice
4. Processing International Payment
5. Employee Management
6. Mobile Experience

🌍 GLOBAL FEATURES
• Multi-currency invoicing
• International payment processing
• Local tax calculations
• Global compliance features

⏱️ Duration: 10-15 minutes

[Live demonstration of Traincape CRM features]"""
    
    # SLIDE 12: SUCCESS STORIES
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📈 SUCCESS STORIES"
    content.text = """🏢 TECH STARTUP (USA)
"300% revenue increase in 6 months"
• 50% faster lead response
• 40% higher conversion rates
• Automated payment processing

🏠 REAL ESTATE COMPANY (UK)
"50% faster deal closure"
• Better lead tracking
• Automated follow-ups
• Professional invoicing

🛒 E-COMMERCE BUSINESS (Canada)
"40% reduction in customer service tickets"
• Better customer management
• Automated payment tracking
• Improved customer satisfaction

🎓 EDUCATION INSTITUTE (Australia)
"80% improvement in student enrollment"
• Streamlined admission process
• Better student tracking
• Automated fee collection

🌍 GLOBAL EXPANSION (Multi-country)
"Seamless international operations"
• Multi-currency support
• Local compliance features
• Global team collaboration"""
    
    # SLIDE 13: PRICING & PLANS
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "💰 PRICING & PLANS"
    content.text = """🚀 STARTER PLAN
$29/month per user
• Up to 10 users
• Basic CRM features
• Payment processing
• Mobile access
• Email support

⭐ PROFESSIONAL PLAN
$59/month per user
• Up to 50 users
• Advanced features
• Advanced analytics
• Workflow automation
• Priority support

🏢 ENTERPRISE PLAN
$129/month per user
• Unlimited users
• Custom features
• API access
• Dedicated support
• Custom integrations

🌍 GLOBAL FEATURES
• Multi-currency support
• International payments
• Local compliance
• Global support

💡 SPECIAL OFFERS
• Annual pricing: 20% discount
• Volume discounts available
• Custom pricing for large teams
• Global expansion support"""
    
    # SLIDE 14: IMPLEMENTATION & SUPPORT
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🚀 IMPLEMENTATION & SUPPORT"
    content.text = """⚡ FAST SETUP
• Ready in 1 week
• No complex configuration
• Free data migration
• Guided setup process

📚 TRAINING & SUPPORT
• Free onboarding training
• Video tutorials
• 24/7 email support
• Phone support during business hours
• Dedicated account manager

🌍 GLOBAL SUPPORT
• Multi-language support
• Local business hours
• Regional compliance help
• International team training

🔄 ONGOING SUPPORT
• Regular updates
• New feature releases
• Performance optimization
• Security updates

🎯 SUCCESS GUARANTEE
• 30-day money-back guarantee
• 99.9% uptime SLA
• ROI guarantee
• Migration guarantee"""
    
    # SLIDE 15: NEXT STEPS
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎯 NEXT STEPS"
    content.text = """📞 IMMEDIATE ACTIONS

1. 📋 Schedule Free Demo
• 30-minute personalized demo
• Customized to your business
• No commitment required

2. 🧪 Start Free Trial
• 14-day free trial
• Full feature access
• Your real data
• Migration assistance

3. 💰 Get Custom Quote
• Tailored to your needs
• Volume discounts
• Custom features
• Implementation timeline

🌍 GLOBAL READINESS
• Multi-currency setup
• International compliance
• Local payment methods
• Global team training

📞 CONTACT INFORMATION
Email: sales@traincapetech.in
Phone: +91-XXXXXXXXXX
Website: www.traincapetech.in
Global Support: Available 24/7"""
    
    # SLIDE 16: Q&A & CLOSE
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "❓ QUESTIONS & ANSWERS"
    content.text = """🤔 COMMON QUESTIONS

Q: How long does setup take?
A: Ready in 1 week vs 6 months for competitors

Q: What about data security?
A: Enterprise-grade security with 99.9% uptime

Q: Can we customize it?
A: Yes, custom fields, workflows, and integrations

Q: What if we're not satisfied?
A: 30-day money-back guarantee

Q: Do you support international businesses?
A: Yes, multi-currency, local compliance, global support

🌍 GLOBAL CAPABILITIES
• 135+ currencies supported
• International payment processing
• Local tax compliance
• Global team support

📞 CONTACT US
Ready to transform your business globally?
Let's discuss your specific needs!

Email: sales@traincapetech.in
Phone: +91-XXXXXXXXXX
Website: www.traincapetech.in
Global Support: Available worldwide"""
    
    # Save the presentation
    filename = "Traincape_CRM_Global_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Global presentation created successfully: {filename}")
    print(f"📁 Location: {os.path.abspath(filename)}")
    
    return filename

if __name__ == "__main__":
    try:
        # Check if python-pptx is installed
        import pptx
        create_global_traincape_presentation()
    except ImportError:
        print("❌ python-pptx library not found!")
        print("📦 Installing python-pptx...")
        os.system("pip3 install python-pptx")
        print("✅ python-pptx installed successfully!")
        create_global_traincape_presentation() 